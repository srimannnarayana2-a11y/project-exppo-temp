"""
Document ingestion — DOCX, PPTX, XLSX → Markdown for KB.

Each format extracts text + structure into Markdown that feeds
directly into the existing chunk.py → embed.py pipeline.
"""

from __future__ import annotations

import logging
from dataclasses import dataclass

logger = logging.getLogger(__name__)


@dataclass
class DocResult:
    """Result of document ingestion."""
    markdown: str = ""
    title: str = ""
    page_count: int = 0
    format: str = ""


# ── DOCX ingestion ──

async def ingest_docx(path: str = "", file_bytes: bytes = b"") -> DocResult:
    """Extract text from DOCX → Markdown."""
    try:
        from docx import Document
    except ImportError:
        logger.error("python-docx not installed")
        return DocResult(markdown="Error: python-docx not installed", format="docx")

    import io

    try:
        if file_bytes:
            doc = Document(io.BytesIO(file_bytes))
        else:
            doc = Document(path)
    except Exception as e:
        return DocResult(markdown=f"Error: {e}", format="docx")

    md_parts = []
    title = ""

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        style = para.style.name.lower() if para.style else ""

        if "heading 1" in style:
            if not title:
                title = text
            md_parts.append(f"# {text}")
        elif "heading 2" in style:
            md_parts.append(f"## {text}")
        elif "heading 3" in style:
            md_parts.append(f"### {text}")
        elif "list" in style:
            md_parts.append(f"- {text}")
        else:
            # Handle bold/italic runs within paragraph
            runs = []
            for run in para.runs:
                t = run.text
                if run.bold:
                    t = f"**{t}**"
                if run.italic:
                    t = f"*{t}*"
                runs.append(t)
            md_parts.append("".join(runs) if runs else text)

    # Extract tables
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)

        if rows:
            header = rows[0]
            md_parts.append("")
            md_parts.append("| " + " | ".join(header) + " |")
            md_parts.append("| " + " | ".join("---" for _ in header) + " |")
            for row in rows[1:]:
                md_parts.append("| " + " | ".join(row) + " |")
            md_parts.append("")

    return DocResult(
        markdown="\n".join(md_parts),
        title=title or "Untitled Document",
        page_count=len(doc.sections),
        format="docx",
    )


# ── PPTX ingestion ──

async def ingest_pptx(path: str = "", file_bytes: bytes = b"") -> DocResult:
    """Extract text from PPTX → Markdown (one section per slide)."""
    try:
        from pptx import Presentation
    except ImportError:
        logger.error("python-pptx not installed")
        return DocResult(markdown="Error: python-pptx not installed", format="pptx")

    import io

    try:
        if file_bytes:
            prs = Presentation(io.BytesIO(file_bytes))
        else:
            prs = Presentation(path)
    except Exception as e:
        return DocResult(markdown=f"Error: {e}", format="pptx")

    md_parts = []
    title = ""

    for i, slide in enumerate(prs.slides):
        slide_title = ""
        slide_content = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        if not slide_title and shape.shape_type == 13:  # Title
                            slide_title = text
                        else:
                            slide_content.append(text)
            elif hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text.strip())

        # Try to get title from first text
        if not slide_title and slide_content:
            slide_title = slide_content.pop(0)

        if not title and slide_title:
            title = slide_title

        md_parts.append(f"\n## Slide {i + 1}: {slide_title or 'Untitled'}\n")
        for line in slide_content:
            md_parts.append(f"- {line}")

        # Extract notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                md_parts.append(f"\n> **Speaker notes:** {notes}")

    return DocResult(
        markdown="\n".join(md_parts),
        title=title or "Untitled Presentation",
        page_count=len(prs.slides),
        format="pptx",
    )


# ── XLSX ingestion ──

async def ingest_xlsx(path: str = "", file_bytes: bytes = b"") -> DocResult:
    """Extract data from XLSX → Markdown tables."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        logger.error("openpyxl not installed")
        return DocResult(markdown="Error: openpyxl not installed", format="xlsx")

    import io

    try:
        if file_bytes:
            wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        else:
            wb = load_workbook(path, read_only=True, data_only=True)
    except Exception as e:
        return DocResult(markdown=f"Error: {e}", format="xlsx")

    md_parts = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        md_parts.append(f"\n## Sheet: {sheet_name}\n")

        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            md_parts.append("*(empty sheet)*")
            continue

        # First row as header
        header = [str(c or "") for c in rows[0]]
        md_parts.append("| " + " | ".join(header) + " |")
        md_parts.append("| " + " | ".join("---" for _ in header) + " |")

        for row in rows[1:100]:  # Cap at 100 rows for sanity
            cells = [str(c or "") for c in row]
            md_parts.append("| " + " | ".join(cells) + " |")

        if len(rows) > 100:
            md_parts.append(f"\n*... and {len(rows) - 100} more rows*")

    wb.close()

    return DocResult(
        markdown="\n".join(md_parts),
        title=wb.sheetnames[0] if wb.sheetnames else "Untitled Spreadsheet",
        page_count=len(wb.sheetnames),
        format="xlsx",
    )


# ── Dispatch by extension ──

DOC_INGESTORS = {
    ".docx": ingest_docx,
    ".pptx": ingest_pptx,
    ".xlsx": ingest_xlsx,
}


async def ingest_document(path: str, file_bytes: bytes = b"") -> DocResult:
    """Auto-detect format and ingest."""
    import os
    ext = os.path.splitext(path)[1].lower()
    ingestor = DOC_INGESTORS.get(ext)
    if not ingestor:
        return DocResult(markdown=f"Unsupported format: {ext}", format=ext)
    return await ingestor(path=path, file_bytes=file_bytes)

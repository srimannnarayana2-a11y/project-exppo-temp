"""
Output renderer — Markdown → PDF / DOCX / PPTX / HTML.

From chats L328: "LLM produces Markdown, deterministic library renders it.
                  PDF proven before DOCX/PPTX. Opt-in per query, not automatic."

Pattern: LLM always outputs Markdown. This module renders that Markdown
into the requested format. Images are embedded when provided.

Dependencies (install only what's needed):
  - PDF:  weasyprint (needs system libs) OR fallback to basic HTML
  - DOCX: python-docx (pure Python)
  - PPTX: python-pptx (pure Python)
  - HTML: built-in (markdown → HTML via basic converter)
"""

from __future__ import annotations

import io
import logging
import re
from typing import Optional

logger = logging.getLogger(__name__)


class RenderResult:
    """Result of a rendering operation."""
    def __init__(
        self,
        content: bytes,
        content_type: str,
        filename: str,
        format_name: str,
    ):
        self.content = content
        self.content_type = content_type
        self.filename = filename
        self.format_name = format_name


# ── Markdown passthrough ──

def render_markdown(md_text: str, title: str = "Report") -> RenderResult:
    """Return Markdown as-is (for download or display)."""
    return RenderResult(
        content=md_text.encode("utf-8"),
        content_type="text/markdown",
        filename=f"{_slugify(title)}.md",
        format_name="markdown",
    )


# ── HTML rendering ──

_HTML_TEMPLATE = """<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title}</title>
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{
            font-family: 'Inter', 'Segoe UI', system-ui, -apple-system, sans-serif;
            line-height: 1.7; color: #1a1a2e; background: #fafbfc;
            max-width: 800px; margin: 0 auto; padding: 2rem;
        }}
        h1 {{ font-size: 2rem; font-weight: 700; margin: 1.5rem 0 1rem; color: #0f0f23; }}
        h2 {{ font-size: 1.5rem; font-weight: 600; margin: 1.3rem 0 0.8rem; color: #16213e; }}
        h3 {{ font-size: 1.2rem; font-weight: 600; margin: 1rem 0 0.6rem; color: #1a1a2e; }}
        p {{ margin: 0.6rem 0; }}
        code {{
            background: #f0f1f3; padding: 0.15rem 0.4rem;
            border-radius: 4px; font-size: 0.9em; font-family: 'JetBrains Mono', monospace;
        }}
        pre {{
            background: #1a1a2e; color: #e8e8e8; padding: 1.2rem;
            border-radius: 8px; overflow-x: auto; margin: 1rem 0;
        }}
        pre code {{ background: none; padding: 0; color: inherit; }}
        ul, ol {{ margin: 0.5rem 0 0.5rem 1.5rem; }}
        li {{ margin: 0.3rem 0; }}
        blockquote {{
            border-left: 4px solid #667eea; padding: 0.5rem 1rem;
            margin: 1rem 0; background: #f8f9ff; border-radius: 0 6px 6px 0;
        }}
        table {{ border-collapse: collapse; width: 100%; margin: 1rem 0; }}
        th, td {{ border: 1px solid #ddd; padding: 0.6rem; text-align: left; }}
        th {{ background: #f0f1f3; font-weight: 600; }}
        img {{ max-width: 100%; border-radius: 8px; margin: 1rem 0; }}
        a {{ color: #667eea; text-decoration: none; }}
        a:hover {{ text-decoration: underline; }}
        hr {{ border: none; border-top: 1px solid #e0e0e0; margin: 1.5rem 0; }}
    </style>
</head>
<body>
{body}
</body>
</html>"""


def render_html(md_text: str, title: str = "Report") -> RenderResult:
    """Convert Markdown to styled HTML."""
    html_body = _md_to_html(md_text)
    full_html = _HTML_TEMPLATE.format(title=title, body=html_body)

    return RenderResult(
        content=full_html.encode("utf-8"),
        content_type="text/html",
        filename=f"{_slugify(title)}.html",
        format_name="html",
    )


# ── PDF rendering ──

def render_pdf(md_text: str, title: str = "Report") -> RenderResult:
    """Convert Markdown → HTML → PDF.

    Tries WeasyPrint first. Falls back to returning HTML with
    a note that PDF rendering needs WeasyPrint installed.
    """
    try:
        from weasyprint import HTML as WPHTML
        html_result = render_html(md_text, title)
        pdf_bytes = WPHTML(string=html_result.content.decode()).write_pdf()
        return RenderResult(
            content=pdf_bytes,
            content_type="application/pdf",
            filename=f"{_slugify(title)}.pdf",
            format_name="pdf",
        )
    except ImportError:
        logger.warning("WeasyPrint not installed, falling back to HTML")
        html_result = render_html(md_text, title)
        return RenderResult(
            content=html_result.content,
            content_type="text/html",
            filename=f"{_slugify(title)}.html",
            format_name="html_fallback",
        )


# ── DOCX rendering ──

def render_docx(md_text: str, title: str = "Report") -> RenderResult:
    """Convert Markdown to DOCX using python-docx."""
    try:
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        logger.warning("python-docx not installed")
        return render_markdown(md_text, title)

    doc = Document()

    # Title
    title_para = doc.add_heading(title, level=0)

    # Parse markdown line by line
    lines = md_text.split("\n")
    in_code_block = False
    code_buffer = []

    for line in lines:
        # Code blocks
        if line.strip().startswith("```"):
            if in_code_block:
                # End code block
                code_text = "\n".join(code_buffer)
                p = doc.add_paragraph()
                run = p.add_run(code_text)
                run.font.name = "Courier New"
                run.font.size = Pt(9)
                code_buffer = []
            in_code_block = not in_code_block
            continue

        if in_code_block:
            code_buffer.append(line)
            continue

        # Headings
        if line.startswith("### "):
            doc.add_heading(line[4:], level=3)
        elif line.startswith("## "):
            doc.add_heading(line[3:], level=2)
        elif line.startswith("# "):
            doc.add_heading(line[2:], level=1)
        elif line.startswith("---"):
            doc.add_paragraph("_" * 40)
        elif line.startswith("- ") or line.startswith("* "):
            doc.add_paragraph(line[2:], style="List Bullet")
        elif re.match(r"^\d+\. ", line):
            text = re.sub(r"^\d+\. ", "", line)
            doc.add_paragraph(text, style="List Number")
        elif line.startswith("> "):
            p = doc.add_paragraph()
            p.style = doc.styles["Quote"] if "Quote" in [s.name for s in doc.styles] else None
            p.add_run(line[2:])
        elif line.strip():
            # Handle inline formatting
            text = line
            p = doc.add_paragraph()
            # Simple bold/italic handling
            parts = re.split(r"(\*\*.*?\*\*|\*.*?\*|`.*?`)", text)
            for part in parts:
                if part.startswith("**") and part.endswith("**"):
                    run = p.add_run(part[2:-2])
                    run.bold = True
                elif part.startswith("*") and part.endswith("*"):
                    run = p.add_run(part[1:-1])
                    run.italic = True
                elif part.startswith("`") and part.endswith("`"):
                    run = p.add_run(part[1:-1])
                    run.font.name = "Courier New"
                    run.font.size = Pt(9)
                else:
                    p.add_run(part)

    # Write to bytes
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)

    return RenderResult(
        content=buf.read(),
        content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        filename=f"{_slugify(title)}.docx",
        format_name="docx",
    )


# ── PPTX rendering ──

def render_pptx(md_text: str, title: str = "Report") -> RenderResult:
    """Convert Markdown to PPTX using python-pptx.

    Strategy: each ## heading becomes a new slide. Content under it
    becomes bullet points on that slide.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        logger.warning("python-pptx not installed")
        return render_markdown(md_text, title)

    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if slide.placeholders[1]:
        slide.placeholders[1].text = "Generated by Agent"

    # Parse sections: each ## becomes a slide
    sections = _split_by_headings(md_text)

    for section_title, section_content in sections:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = section_title

        # Add content as bullet points
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()

        lines = section_content.strip().split("\n")
        first = True
        for line in lines:
            if not line.strip():
                continue
            # Clean markdown formatting
            clean = re.sub(r"\*\*(.*?)\*\*", r"\1", line)
            clean = re.sub(r"\*(.*?)\*", r"\1", clean)
            clean = re.sub(r"`(.*?)`", r"\1", clean)
            clean = clean.lstrip("- *•")

            if first:
                tf.paragraphs[0].text = clean
                first = False
            else:
                p = tf.add_paragraph()
                p.text = clean
                p.font.size = Pt(14)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    return RenderResult(
        content=buf.read(),
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=f"{_slugify(title)}.pptx",
        format_name="pptx",
    )


# ── Dispatch ──

RENDERERS = {
    "markdown": render_markdown,
    "md": render_markdown,
    "html": render_html,
    "pdf": render_pdf,
    "docx": render_docx,
    "pptx": render_pptx,
}


def render_output(
    md_text: str,
    format: str = "markdown",
    title: str = "Report",
) -> RenderResult:
    """Render markdown to the requested format."""
    renderer = RENDERERS.get(format.lower(), render_markdown)
    return renderer(md_text, title)


# ── Helpers ──

def _slugify(text: str) -> str:
    """Convert text to filename-safe slug."""
    return re.sub(r"[^\w\-]", "_", text.lower())[:50]


def _md_to_html(md_text: str) -> str:
    """Basic Markdown → HTML converter (no external deps)."""
    lines = md_text.split("\n")
    html_lines = []
    in_code = False
    in_list = False
    in_table = False

    for line in lines:
        # Code blocks
        if line.strip().startswith("```"):
            if in_code:
                html_lines.append("</code></pre>")
            else:
                lang = line.strip()[3:]
                html_lines.append(f'<pre><code class="language-{lang}">')
            in_code = not in_code
            continue

        if in_code:
            html_lines.append(_escape_html(line))
            continue

        # Tables
        if "|" in line and line.strip().startswith("|"):
            cells = [c.strip() for c in line.strip().strip("|").split("|")]
            if all(re.match(r"^[-:]+$", c) for c in cells):
                continue  # separator row
            if not in_table:
                html_lines.append("<table>")
                in_table = True
            tag = "th" if not in_table or html_lines[-1] == "<table>" else "td"
            row = "".join(f"<{tag}>{_inline_md(c)}</{tag}>" for c in cells)
            html_lines.append(f"<tr>{row}</tr>")
            continue
        elif in_table:
            html_lines.append("</table>")
            in_table = False

        # Headings
        if line.startswith("### "):
            html_lines.append(f"<h3>{_inline_md(line[4:])}</h3>")
        elif line.startswith("## "):
            html_lines.append(f"<h2>{_inline_md(line[3:])}</h2>")
        elif line.startswith("# "):
            html_lines.append(f"<h1>{_inline_md(line[2:])}</h1>")
        elif line.startswith("---"):
            html_lines.append("<hr>")
        elif line.startswith("> "):
            html_lines.append(f"<blockquote><p>{_inline_md(line[2:])}</p></blockquote>")
        elif line.startswith("- ") or line.startswith("* "):
            if not in_list:
                html_lines.append("<ul>")
                in_list = True
            html_lines.append(f"<li>{_inline_md(line[2:])}</li>")
        else:
            if in_list:
                html_lines.append("</ul>")
                in_list = False
            if line.strip():
                html_lines.append(f"<p>{_inline_md(line)}</p>")

    if in_list:
        html_lines.append("</ul>")
    if in_table:
        html_lines.append("</table>")

    return "\n".join(html_lines)


def _inline_md(text: str) -> str:
    """Convert inline markdown (bold, italic, code, links)."""
    text = re.sub(r"\*\*(.*?)\*\*", r"<strong>\1</strong>", text)
    text = re.sub(r"\*(.*?)\*", r"<em>\1</em>", text)
    text = re.sub(r"`(.*?)`", r"<code>\1</code>", text)
    text = re.sub(r"\[(.*?)\]\((.*?)\)", r'<a href="\2">\1</a>', text)
    text = re.sub(r"!\[(.*?)\]\((.*?)\)", r'<img src="\2" alt="\1">', text)
    return text


def _escape_html(text: str) -> str:
    """Escape HTML special chars."""
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def _split_by_headings(md_text: str) -> list[tuple[str, str]]:
    """Split markdown by ## headings into (title, content) pairs."""
    sections = []
    current_title = "Overview"
    current_content = []

    for line in md_text.split("\n"):
        if line.startswith("## "):
            if current_content:
                sections.append((current_title, "\n".join(current_content)))
            current_title = line[3:]
            current_content = []
        elif line.startswith("# "):
            # Skip top-level heading (used as title slide)
            continue
        else:
            current_content.append(line)

    if current_content:
        sections.append((current_title, "\n".join(current_content)))

    return sections
